"""
Create PowerPoint presentation for Taobao Maternity Analysis
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Helper functions
def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title, content_list):
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    
    for i, text in enumerate(content_list):
        if i == 0:
            tf.text = text
        else:
            p = tf.add_paragraph()
            p.text = text
            p.level = 0
    return slide

# Slide 1: Title
add_title_slide(prs, 
    "淘宝母婴购物数据分析",
    "Taobao Maternity Shopping Data Analysis\\n\\n基于阿里云天池数据集\\n2026年2月")

# Slide 2: Project Overview
add_content_slide(prs, "项目概述", [
    "分析目标：洞察母婴电商平台用户行为和商品偏好",
    "数据规模：1,000位用户，15,000条交易记录",
    "时间跨度：2012-2015年",
    "核心指标：用户画像、商品分析、时间趋势、RFM价值分析",
    "数据来源：阿里云天池数据集 #45"
])

# Slide 3: Key Findings
add_content_slide(prs, "核心发现", [
    "冠军客户占比14.5%，是平台核心收入来源",
    "18.7%客户存在流失风险，需立即召回",
    "1-2岁是购买高峰期，占比49.4%",
    "品类销量相对均衡，差异<5%",
    "平台年增长率稳定在2-3%"
])

# Slide 4: User Profile
add_content_slide(prs, "用户画像分析", [
    "性别分布：男女婴比例均衡（46% vs 44%）",
    "黄金年龄段：1-2岁是核心购买期",
    "年龄与频次：3岁以上用户购买频次最高（6.45次/人）",
    "增长潜力：0-12个月是潜力培育期"
])

# Slide 5: Product Analysis
add_content_slide(prs, "商品品类分析", [
    "品类均衡：Top 5品类销量差异<5%",
    "广泛覆盖：各品类覆盖约95%用户",
    "订单结构：85%订单≤5件，中小订单为主",
    "热销品类：50014815略领先，但优势不明显"
])

# Slide 6: Time Trend
add_content_slide(prs, "时间趋势分析", [
    "年度增长：年增长率2-3%，增长稳定",
    "用户粘性：用户基数稳定在970-980人",
    "无季节性：全年销售均匀分布",
    "复购率高：平台处于稳定增长期"
])

# Slide 7: RFM Analysis
add_content_slide(prs, "RFM客户价值分析", [
    "冠军客户（14.5%）：高活跃+高频次+高消费",
    "重要保持客户（18.0%）：购买频次稳定",
    "重要挽留客户（18.7%）：曾经活跃但近期未购买",
    "重要发展客户（11.3%）：新客户，潜力待挖掘",
    "不可失去客户（5.9%）：曾经高价值，近期沉默"
])

# Slide 8: Business Recommendations
add_content_slide(prs, "业务建议", [
    "P0 - 冠军客户VIP专属服务：145位核心客户",
    "P0 - 流失客户召回活动：187位风险客户",
    "P1 - 忠诚客户升单计划：180位忠诚客户",
    "P1 - 新客户培育计划：113位新客户",
    "P2 - 1-2岁用户专项运营：核心年龄段"
])

# Slide 9: Thank you
add_title_slide(prs, "谢谢观看", "Thank You\\n\\n数据分析驱动业务增长")

# Save presentation
output_path = "Taobao_Maternity_Analysis_Report.pptx"
prs.save(output_path)
print(f"PPT报告已生成: {output_path}")
